"""文档解析：把 PDF / PPTX / DOCX / TXT / Markdown 提取为 (位置, 文本) 片段列表。"""
import os

SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".txt", ".md", ".markdown"}


def extract_segments(path, filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return _pdf(path)
    if ext == ".pptx":
        return _pptx(path)
    if ext == ".docx":
        return _docx(path)
    if ext in (".txt", ".md", ".markdown"):
        return [("", _read_text(path))]
    if ext in (".ppt", ".doc"):
        raise ValueError(
            "暂不支持旧版 {0} 格式，请在 Office / WPS 中另存为 {0}x 格式后再上传".format(ext)
        )
    raise ValueError(
        "暂不支持 {} 格式，目前支持 PDF、PPTX、DOCX、TXT、Markdown".format(ext or "未知")
    )


def _pdf(path):
    from pypdf import PdfReader

    reader = PdfReader(path)
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            raise ValueError("PDF 已加密，无法解析，请先解除密码保护")
    segs = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            segs.append(("第 {} 页".format(i + 1), text))
    if not segs:
        raise ValueError("未能从 PDF 中提取到文字（可能是扫描版 / 纯图片 PDF）")
    return segs


def _shape_texts(shape, out):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _shape_texts(sub, out)
        return
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        out.append(shape.text_frame.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                out.append(" | ".join(cells))


def _pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    segs = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            _shape_texts(shape, texts)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append("【演讲者备注】" + notes)
        body = "\n".join(texts).strip()
        if body:
            segs.append(("第 {} 页幻灯片".format(i + 1), body))
    if not segs:
        raise ValueError("未能从 PPTX 中提取到文字内容")
    return segs


def _docx(path):
    from docx import Document

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    body = "\n".join(parts).strip()
    if not body:
        raise ValueError("未能从 DOCX 中提取到文字内容")
    return [("", body)]


def _read_text(path):
    with open(path, "rb") as f:
        raw = f.read()
    for enc in ("utf-8-sig", "utf-8", "gb18030", "big5", "latin-1"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError("无法识别文件编码")
